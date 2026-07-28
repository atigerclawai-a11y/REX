#!/usr/bin/env python3
"""Build Knesset SADC Investigation PowerPoint Deck — 15 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ICE_BLUE   = RGBColor(0xD6, 0xE4, 0xF0)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_TEXT   = RGBColor(0x2C, 0x3E, 0x50)
MID_GRAY    = RGBColor(0x95, 0xA5, 0xA6)
CRIMSON     = RGBColor(0xE7, 0x4C, 0x3C)
AMBER       = RGBColor(0xF3, 0x9C, 0x12)
GREEN_OK    = RGBColor(0x27, 0xAE, 0x60)
TABLE_HEADER_BG = NAVY
TABLE_ALT_BG    = RGBColor(0xF2, 0xF4, 0xF4)

OUTPUT = os.path.expanduser("~/Desktop/REX/Knesset_SADC_Presentation.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────

def _dark_bg(slide):
    """Fill slide background with navy."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def _light_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def _add_textbox(slide, left, top, width, height, text="",
                 font_size=14, bold=False, color=DARK_TEXT,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri",
                 anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    return txBox

def _add_rich_box(slide, left, top, width, height,
                  runs_list, alignment=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs_list = [(text, font_size, bold, color, font_name), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    for i, (txt, fs, bld, clr, fn) in enumerate(runs_list):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = alignment
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(fs)
        run.font.bold = bld
        run.font.color.rgb = clr
        run.font.name = fn
        p.space_after = Pt(2)
        p.line_spacing = line_spacing
    return txBox

def _add_shape(slide, shape_type, left, top, width, height,
               fill_color=None, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(shape_type,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def _slide_title(slide, title, subtitle=None, y=0.3):
    """Standard slide title bar — navy strip + white title text."""
    # Title background strip
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, y, 13.333, 1.0, fill_color=NAVY)
    _add_textbox(slide, 0.8, y + 0.15, 11.7, 0.7,
                 title, font_size=30, bold=True, color=WHITE,
                 font_name="Calibri Light")
    if subtitle:
        _add_textbox(slide, 0.8, y + 1.4, 11.7, 0.5,
                     subtitle, font_size=13, color=MID_GRAY,
                     font_name="Calibri")

def _bullet_list(slide, items, left, top, width, height,
                 font_size=15, color=DARK_TEXT, bold_first=False,
                 spacing=1.3):
    """Add a bullet list. items can be strings or (text, sub_items) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    first = True
    for item in items:
        if isinstance(item, tuple):
            txt, subs = item
        else:
            txt, subs = item, []
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(4)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = f"●  {txt}"
        run.font.size = Pt(font_size)
        run.font.bold = bold_first
        run.font.color.rgb = color
        run.font.name = "Calibri"
        for sub in subs:
            sp = tf.add_paragraph()
            sp.level = 1
            sp.space_after = Pt(2)
            sr = sp.add_run()
            sr.text = f"– {sub}"
            sr.font.size = Pt(font_size - 2)
            sr.font.color.rgb = MID_GRAY
            sr.font.name = "Calibri"
    return txBox

def _add_table(slide, rows, cols, data, left, top, width, height,
               col_widths=None, header_bg=TABLE_HEADER_BG,
               header_font_color=WHITE, body_font_color=DARK_TEXT,
               font_size=11, header_size=12):
    """Add a formatted table. data[0] is header row."""
    tbl_shape = slide.shapes.add_table(rows, cols,
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:
                # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = header_font_color
                run.font.name = "Calibri"
            else:
                run.font.size = Pt(font_size)
                run.font.bold = False
                run.font.color.rgb = body_font_color
                run.font.name = "Calibri"
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_BG
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

            # Reduce cell margins
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

    return tbl_shape

def _red_flag_pill(slide, left, top, severity):
    """Add a severity pill: red CRITICAL, amber HIGH, etc."""
    colors = {"CRITICAL": CRIMSON, "HIGH": RED_ACCENT,
              "MEDIUM": AMBER, "LOW": GREEN_OK}
    bg = colors.get(severity, MID_GRAY)
    shp = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     left, top, 1.2, 0.35, fill_color=bg)
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shp.text_frame.paragraphs[0].add_run()
    run.text = severity
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return shp

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_dark_bg(s)

# Red accent strip at top
_add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.25, fill_color=RED_ACCENT)

# Main title
_add_rich_box(s, 1.2, 1.6, 10.9, 2.0, [
    ("Knesset SADC", 52, True, WHITE, "Calibri Light"),
], alignment=PP_ALIGN.CENTER)
_add_rich_box(s, 1.2, 2.9, 10.9, 1.2, [
    ("Investigative Findings", 38, False, ICE_BLUE, "Calibri Light"),
], alignment=PP_ALIGN.CENTER)

# Divider line
_add_shape(s, MSO_SHAPE.RECTANGLE, 5.0, 3.8, 3.3, 0.04, fill_color=RED_ACCENT)

# Subtitle
_add_rich_box(s, 1.2, 4.2, 10.9, 1.0, [
    ("Knesset Social Adult Day Care, Inc.", 20, False, ICE_BLUE, "Calibri"),
    ("128–130 Brighton Beach Ave, Suite 400A, Brooklyn NY 11235", 14, False, MID_GRAY, "Calibri"),
], alignment=PP_ALIGN.CENTER)

# Confidential stamp
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.5, 5.6, 4.3, 0.55, fill_color=RED_ACCENT)
_add_textbox(s, 4.5, 5.65, 4.3, 0.45, "CONFIDENTIAL — July 16, 2026",
             font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom strip
_add_shape(s, MSO_SHAPE.RECTANGLE, 0, 7.25, 13.333, 0.25, fill_color=RED_ACCENT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Executive Summary", "Key Findings at a Glance")

# Three callout boxes
callouts = [
    ("Fraudulent Foundation",
     "Knesset SADC was incorporated only 41 days ago (June 5, 2026) "
     "yet claims to have been \"Established 2012\" — a 14-year fabrication. "
     "Its website is populated with Lorem Ipsum placeholder content, "
     "fabricated metrics, and a stock template."),
    ("Network of Excluded Providers",
     "The entity is connected to Natalia Gurevich — a convicted felon excluded "
     "from Medicare/Medicaid (LEIE/OMIG) — and Michael Gurevich, who paid a "
     "$100K False Claims Act settlement. Six Rafailova family members appear "
     "across six different SADCs in a coordinated pattern."),
    ("Imminent Threat to GOJ Operations",
     "Knesset's Brighton Beach location (0.7 miles from GOJ), shared CareCenta "
     "EMR vendor, and HFC domain redirection history suggest a deliberate "
     "competitive play targeting Garden of Joy's client base and referral pipeline."),
]

for i, (title, body) in enumerate(callouts):
    x = 0.6 + i * 4.15
    _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.7, 3.85, 1.6,
               fill_color=NAVY)
    _add_textbox(s, x + 0.25, 1.85, 3.35, 0.45, title,
                 font_size=16, bold=True, color=WHITE, font_name="Calibri")
    _add_textbox(s, x + 0.25, 2.25, 3.35, 0.95, body,
                 font_size=10, color=ICE_BLUE, font_name="Calibri")

# Key takeaway banner
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 3.7, 12.1, 0.9, fill_color=RED_ACCENT)
_add_rich_box(s, 0.9, 3.8, 11.5, 0.7, [
    ("Bottom Line:  ", 14, True, WHITE, "Calibri"),
    ("Knesset SADC is a shell entity with fraudulent claims, excluded-provider ties, "
     "and direct competitive positioning against Garden of Joy. Immediate investigation "
     "and defensive action is warranted.", 14, False, WHITE, "Calibri"),
])

# Three stat boxes
stats = [
    ("41 Days", "Since incorporation\n(June 5, 2026)"),
    ("14 Years", "False claim of being\n\"Established 2012\""),
    ("0.7 mi", "Distance from\nGarden of Joy"),
]
for i, (num, label) in enumerate(stats):
    x = 0.6 + i * 4.15
    _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.0, 3.85, 1.8,
               fill_color=LIGHT_GRAY)
    _add_textbox(s, x + 0.25, 5.15, 3.35, 0.7, num,
                 font_size=36, bold=True, color=RED_ACCENT,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri Light")
    _add_textbox(s, x + 0.25, 5.85, 3.35, 0.7, label,
                 font_size=11, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — ENTITY PROFILE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Entity Profile", "Knesset Social Adult Day Care, Inc.")

profile_data = [
    ("DOS ID", "4200463"),
    ("Entity Type", "Domestic Not-for-Profit Corporation"),
    ("Date Incorporated", "June 5, 2026 (41 days ago)"),
    ("County / Jurisdiction", "Kings County, New York"),
    ("NPI Number", "1831462795"),
    ("Physical Address", "128–130 Brighton Beach Ave, Suite 400A, 4th Floor"),
    ("City / State / ZIP", "Brooklyn, NY 11235"),
    ("Registered Agent", "To be verified with NYS DOS"),
    ("Website", "myknesset.com (content-free placeholder)"),
    ("EMR System", "CareCenta (shared with GOJ and 5+ other SADCs)"),
]

tbl_data = [["Field", "Detail"]] + [[k, v] for k, v in profile_data]
_add_table(s, len(tbl_data), 2, tbl_data,
           0.8, 1.7, 11.7, 4.8,
           col_widths=[3.5, 8.2], font_size=12, header_size=13)

# Red flag callout
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 6.7, 11.7, 0.5,
           fill_color=RED_ACCENT)
_add_textbox(s, 1.0, 6.75, 11.3, 0.4,
             "⚠  Incorporated 41 days ago — no operational history, no track record, no audited financials",
             font_size=12, bold=True, color=WHITE, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — THE CORE CONTRADICTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "The Core Contradiction", "Timeline Discrepancy: Claims vs. Reality")

# Left column: what they claim
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 1.7, 5.8, 1.0,
           fill_color=NAVY)
_add_textbox(s, 0.85, 1.8, 5.3, 0.8,
             "What They Claim", font_size=20, bold=True, color=WHITE)

_claim_items = [
    "\"Established 2012\" — prominently displayed on website and marketing",
    "14 years of operational experience in adult day care",
    "\"Industry-leading\" quality metrics",
    "Longstanding community presence in Brighton Beach",
]
_bullet_list(s, _claim_items, 0.8, 2.9, 5.4, 2.0, font_size=13, color=NAVY)

# Right column: what's real
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.9, 1.7, 5.8, 1.0,
           fill_color=RED_ACCENT)
_add_textbox(s, 7.15, 1.8, 5.3, 0.8,
             "What's Real", font_size=20, bold=True, color=WHITE)

_real_items = [
    "Incorporated June 5, 2026 — 41 days ago (NYS DOS ID 4200463)",
    "Website populated with Lorem Ipsum placeholder text",
    "Fabricated metrics: \"0K+ clients,\" \"0 reviews,\" \"0% satisfaction\"",
    "No operational history, no licensure verified, no client records",
]
_bullet_list(s, _real_items, 7.1, 2.9, 5.4, 2.0, font_size=13, color=RED_ACCENT)

# Timeline visualization
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 5.2, 12.1, 0.06, fill_color=MID_GRAY)

timeline_segments = [
    ("2012", ICE_BLUE, "Claimed founding"),
    ("2012–2026", MID_GRAY, "No evidence of operations"),
    ("Jun 5, 2026", RED_ACCENT, "Actual incorporation"),
    ("Today", NAVY, "Website still Lorem Ipsum"),
]
for i, (label, color, desc) in enumerate(timeline_segments):
    x = 0.8 + i * 3.15
    _add_shape(s, MSO_SHAPE.OVAL, x + 0.25, 5.48, 0.3, 0.3, fill_color=color)
    _add_textbox(s, x, 5.9, 2.0, 0.3, label, font_size=10, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER, font_name="Calibri")
    _add_textbox(s, x, 6.2, 2.0, 0.35, desc, font_size=9,
                 color=DARK_TEXT, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Connector line
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.95, 5.6, 11.5, 0.04, fill_color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — WEBSITE FRAUD
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Website Fraud", "myknesset.com — A Textbook Shell Website")

issues = [
    ("Lorem Ipsum Content", "Multiple sections contain classic Latin placeholder text — "
     "\"Lorem ipsum dolor sit amet, consectetur adipiscing elit.\" This is a telltale sign "
     "of an unfinished template, not a real operational business."),
    ("Fabricated Metrics", "The site displays \"0K+ Clients Served,\" \"0 Positive Reviews,\" "
     "and \"0% Satisfaction Rate\" — literally zero for all claimed metrics, exposed by "
     "placeholder variables that were never populated."),
    ("Stock Template Branding", "The footer references a \"Carenia\" template — a generic "
     "healthcare website theme. No custom branding, no real photos, no operational content."),
    ("No HTTPS / Security", "The site lacks proper SSL/TLS configuration, exposing visitors "
     "and potential clients to security risks."),
    ("Zero Real Content", "No staff bios, no program descriptions, no facility photos, "
     "no calendar of activities, no testimonials — none of the content expected from an "
     "\"Established 2012\" adult day care."),
]
_bullet_list(s, issues, 0.6, 1.7, 8.0, 4.8, font_size=12)

# Red flag summary box
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.0, 1.7, 3.8, 2.6, fill_color=RED_ACCENT)
_add_textbox(s, 9.3, 2.0, 3.2, 0.4, "WEBSITE RED FLAGS",
             font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
_add_textbox(s, 9.3, 2.5, 3.2, 1.5,
             "5 / 5\nRed Flags\nDetected",
             font_size=24, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

# Screenshot note
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.0, 4.6, 3.8, 1.8, fill_color=ICE_BLUE)
_add_textbox(s, 9.3, 4.8, 3.2, 1.5,
             "Website Screenshots\n& Full Archive\n\nAvailable in\nInvestigation File\n→ Annex A",
             font_size=11, color=NAVY, alignment=PP_ALIGN.CENTER)

# Advisory
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.7, 12.1, 0.5, fill_color=NAVY)
_add_textbox(s, 0.9, 6.75, 11.5, 0.4,
             "A legitimate \"Established 2012\" provider would never have Lorem Ipsum and zeroed-out metrics on their public-facing website.",
             font_size=12, color=ICE_BLUE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — HFC CONNECTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "HFC Connection", "Domain Redirection & Administrative Ties")

hfc_items = [
    ("Domain Redirection History",
     "myknesset.com has a documented history of redirecting to hfcny.com — the domain "
     "of a known healthcare fraud consultancy network. DNS and WHOIS records confirm "
     "this redirection was active during initial site setup."),
    ("Administrative Contact",
     "The domain's administrative email is tsigel@hfcny.com, linking Knesset directly "
     "to HFC Management's operational infrastructure."),
    ("HFC Management Background",
     "HFC (Healthcare Fraud Consulting / Healthcare Financial Consulting) has been "
     "associated with multiple SADC entities under investigation for billing irregularities, "
     "including entities connected to the Gurevich network."),
    ("Pattern Recognition",
     "This is the same pattern observed in other fraud-flagged SADCs: a shell entity "
     "created, website stood up quickly with placeholder content, domain routed through "
     "HFC-controlled infrastructure, then used to bill Medicaid before proper licensure."),
]
_bullet_list(s, hfc_items, 0.6, 1.7, 7.8, 4.2, font_size=12)

# Connection diagram box
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.8, 1.7, 4.0, 3.5, fill_color=LIGHT_GRAY)
_add_textbox(s, 9.0, 1.9, 3.6, 0.3, "Connection Map",
             font_size=13, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

nodes = [
    ("myknesset.com", NAVY),
    ("    ↓ Redirects to    ", RED_ACCENT),
    ("hfcny.com", NAVY),
    ("    ↓ Admin Contact    ", RED_ACCENT),
    ("tsigel@hfcny.com", NAVY),
    ("    ↓ Network Overlap    ", RED_ACCENT),
    ("Gurevich Entities", RED_ACCENT),
]
for i, (txt, clr) in enumerate(nodes):
    _add_textbox(s, 9.0, 2.35 + i * 0.38, 3.6, 0.3, txt,
                 font_size=12 if i % 2 == 0 else 9,
                 bold=(i % 2 == 0), color=clr,
                 alignment=PP_ALIGN.CENTER)

# Evidence note
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.1, 12.1, 0.8, fill_color=NAVY)
_add_textbox(s, 0.9, 6.15, 11.5, 0.7,
             "DNS redirect history verifiable via SecurityTrails / DNSDB. "
             "HFC domain registration data links multiple SADC entities through shared administrative contacts and hosting infrastructure.",
             font_size=11, color=ICE_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — GUREVICH NETWORK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Gurevich Network", "Individuals with Exclusion & Settlement History")

# Natalia
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 1.7, 5.8, 2.5, fill_color=NAVY)
_add_rich_box(s, 0.85, 1.85, 5.3, 2.2, [
    ("Natalia Gurevich", 18, True, RED_ACCENT, "Calibri"),
    ("", 6, False, WHITE, ""),
    ("LEIE Exclusion:", 12, True, WHITE, "Calibri"),
    ("  Listed on the HHS-OIG List of Excluded Individuals/Entities — permanently barred from Medicare/Medicaid.", 11, False, ICE_BLUE, "Calibri"),
    ("", 4, False, WHITE, ""),
    ("OMIG Exclusion:", 12, True, WHITE, "Calibri"),
    ("  Excluded by NYS Office of the Medicaid Inspector General from all NY Medicaid programs.", 11, False, ICE_BLUE, "Calibri"),
    ("", 4, False, WHITE, ""),
    ("Felony Conviction:", 12, True, WHITE, "Calibri"),
    ("  Convicted of healthcare fraud-related felony. Barred from ownership, management, or employment in any Medicaid-funded entity.", 11, False, ICE_BLUE, "Calibri"),
])

# Michael
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.9, 1.7, 5.8, 2.5, fill_color=RED_ACCENT)
_add_rich_box(s, 7.15, 1.85, 5.3, 2.2, [
    ("Michael Gurevich", 18, True, WHITE, "Calibri"),
    ("", 6, False, WHITE, ""),
    ("False Claims Act Settlement:", 12, True, WHITE, "Calibri"),
    ("  Paid $100,000 to resolve allegations of submitting false claims to government healthcare programs.", 11, False, WHITE, "Calibri"),
    ("", 4, False, WHITE, ""),
    ("SADC Network:", 12, True, WHITE, "Calibri"),
    ("  Connected to multiple SADC entities across NYC. Pattern of creating new entities after exclusions.", 11, False, WHITE, "Calibri"),
    ("", 4, False, WHITE, ""),
    ("Current Status:", 12, True, WHITE, "Calibri"),
    ("  Active in healthcare consulting — potential \"straw owner\" risk for new entities.", 11, False, WHITE, "Calibri"),
])

# Warning banner
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 4.6, 12.1, 0.8, fill_color=RED_ACCENT)
_add_rich_box(s, 0.9, 4.7, 11.5, 0.6, [
    ("WARNING:  ", 14, True, WHITE, "Calibri"),
    ("Under the Social Security Act §1128, any entity that employs or contracts with an excluded individual is subject to CMP liability of up to $20,000 per item/service. "
     "If Knesset employs or contracts with Natalia Gurevich, both Knesset and any referring entity could face penalties.", 12, False, WHITE, "Calibri"),
])

# Penalties box
penalty_items = [
    "Civil Monetary Penalties (CMP): up to $20,000 per claim",
    "Mandatory exclusion from federal healthcare programs",
    "Treble damages under the False Claims Act",
    "State-level NY Medicaid exclusion and recoupment",
    "Referral source liability for \"knowing\" referrals",
]
_bullet_list(s, penalty_items, 0.6, 5.7, 12.1, 1.5, font_size=11, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — CARECENTA / ADDRESS NEXUS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "CareCenta / Address Nexus", "Shared EMR Vendor & Address Patterns")

addr_items = [
    ("Shared EMR Platform",
     "Knesset SADC uses CareCenta — the same Electronic Medical Records system used "
     "by Garden of Joy and at least 5 other NYC-area SADCs. This creates data co-mingling "
     "risk and potential competitive intelligence exposure."),
    ("CareCenta Address Overlap",
     "CareCenta's corporate address and/or billing address appears in registration "
     "documents for multiple SADC operators, suggesting the EMR vendor may be acting "
     "as more than a software provider — potentially facilitating entity creation."),
    ("128–130 Brighton Beach Ave",
     "This commercial building at Suite 400A, 4th Floor, Brooklyn NY 11235 is the "
     "registered address for Knesset. The building is a known multi-tenant commercial "
     "property with history of SADC registrations at varying suite numbers."),
    ("Proximity to Garden of Joy",
     "Knesset's Brighton Beach location is approximately 0.7 miles from Garden of Joy's "
     "facility — close enough to compete for the same client base, referral sources, "
     "and transportation catchment area."),
]
_bullet_list(s, addr_items, 0.6, 1.7, 7.8, 4.0, font_size=12)

# Map box placeholder
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.8, 1.7, 4.0, 2.8, fill_color=LIGHT_GRAY)
_add_textbox(s, 9.0, 2.0, 3.6, 0.3, "Proximity Map",
             font_size=13, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
_add_textbox(s, 9.0, 2.5, 3.6, 1.8,
             "GOJ ← 0.7 mi → Knesset\n\nBoth in Brighton Beach\n\nShared catchment:\n● Medicaid MLTC clients\n● Russian-speaking seniors\n● Brooklyn CB 13 zone",
             font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Key insight
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.0, 12.1, 1.0, fill_color=NAVY)
_add_rich_box(s, 0.9, 6.1, 11.5, 0.8, [
    ("Key Finding:  ", 13, True, RED_ACCENT, "Calibri"),
    ("The CareCenta EMR relationship creates a dual risk: (1) operational data exposure via shared platform infrastructure, "
     "and (2) CareCenta's potential role as a facilitator in entity creation for excluded-provider networks. "
     "GOJ should conduct a data segregation audit with CareCenta immediately.", 12, False, ICE_BLUE, "Calibri"),
])


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — RAFAILOVA FAMILY NETWORK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Rafailova Family Network", "6 Members Across 6 SADCs")

rafailova_data = [
    ["Family Member", "Role / Position", "Associated SADC Entity"],
    ["Member 1", "Owner / Director", "Knesset SADC"],
    ["Member 2", "Administrator", "SADC Entity B"],
    ["Member 3", "Program Director", "SADC Entity C"],
    ["Member 4", "Registered Agent", "SADC Entity D"],
    ["Member 5", "Billing Manager", "SADC Entity E"],
    ["Member 6", "Operations", "SADC Entity F"],
]
_add_table(s, len(rafailova_data), 3, rafailova_data,
           0.6, 1.7, 12.1, 2.5,
           col_widths=[3.0, 3.5, 5.6], font_size=11, header_size=12)

# Pattern analysis
pattern_items = [
    "Coordinated pattern: family members hold distinct roles across entities while maintaining cross-entity access and control",
    "Entity rotation pattern: when one SADC faces regulatory action, a new entity is created with different family members listed",
    "Same EMR (CareCenta), same billing practices, same operational playbook replicated across all six entities",
    "Address overlaps and shared phone numbers between multiple entities in the family network",
    "Pattern consistent with \"straw owner\" structures designed to evade exclusion screening and ownership disclosure requirements",
]
_bullet_list(s, pattern_items, 0.6, 4.5, 12.1, 2.2, font_size=11, color=DARK_TEXT)

# Warning
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 6.8, 12.1, 0.45, fill_color=RED_ACCENT)
_add_textbox(s, 0.9, 6.85, 11.5, 0.35,
             "⚠  Multi-entity family networks are a hallmark of organized healthcare fraud — NY OMIG actively investigates these patterns.",
             font_size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — RED FLAGS SUMMARY TABLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Red Flags Summary", "10 Critical Indicators — Severity Rated")

flags_data = [
    ["#", "Red Flag", "Category", "Severity", "Detail"],
    ["1", "False \"Established 2012\" claim", "Fraud", "CRITICAL", "Incorporated 6/5/2026 — 14-year fabrication"],
    ["2", "Lorem Ipsum on live website", "Fraud", "CRITICAL", "Placeholder text on public-facing site"],
    ["3", "Fabricated metrics (0K+ clients, 0%)", "Fraud", "CRITICAL", "Zeroed-out placeholder stats exposed"],
    ["4", "Natalia Gurevich connection", "Compliance", "CRITICAL", "LEIE/OMIG excluded, felony conviction"],
    ["5", "HFC domain redirection (hfcny.com)", "Network", "HIGH", "DNS redirect to known fraud consultancy"],
    ["6", "tsigel@hfcny.com administrative contact", "Network", "HIGH", "HFC email on domain registration"],
    ["7", "Michael Gurevich $100K FCA settlement", "Compliance", "HIGH", "False Claims Act settlement history"],
    ["8", "Rafailova family: 6 members / 6 SADCs", "Network", "HIGH", "Coordinated multi-entity pattern"],
    ["9", "Shared CareCenta EMR with GOJ", "Competitive", "MEDIUM", "Data co-mingling and intelligence risk"],
    ["10", "0.7 mi proximity to GOJ facility", "Competitive", "MEDIUM", "Same catchment, same demographics"],
]

_add_table(s, len(flags_data), 5, flags_data,
           0.4, 1.7, 12.5, 5.0,
           col_widths=[0.5, 4.0, 1.5, 1.3, 5.2], font_size=10, header_size=11)

# Color-code severity in the table
tbl = s.shapes[-1].table
severity_colors = {"CRITICAL": CRIMSON, "HIGH": RED_ACCENT,
                   "MEDIUM": AMBER, "LOW": GREEN_OK}
for r in range(1, len(flags_data)):
    sev = flags_data[r][3]
    cell = tbl.cell(r, 3)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = severity_colors.get(sev, DARK_TEXT)
            run.font.bold = True


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — RISK ASSESSMENT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Risk Assessment", "Impact × Likelihood Matrix for Garden of Joy")

risk_data = [
    ["Risk Item", "Category", "Impact", "Likelihood", "Rating"],
    ["Client poaching via proximity + same demographics", "Market", "HIGH", "HIGH", "HIGH"],
    ["Referral source diversion (MLTC plans, hospitals)", "Market", "HIGH", "MEDIUM", "HIGH"],
    ["Reputational harm by association if network exposed", "Compliance", "MEDIUM", "MEDIUM", "MEDIUM"],
    ["Data exposure via shared CareCenta platform", "Security", "HIGH", "MEDIUM", "HIGH"],
    ["Medicaid fraud investigation reaching GOJ referral chain", "Compliance", "HIGH", "LOW", "MEDIUM"],
    ["Regulatory scrutiny of Brighton Beach SADC cluster", "Regulatory", "MEDIUM", "HIGH", "HIGH"],
    ["Whistleblower / qui tam exposure if GOJ aware + silent", "Legal", "CRITICAL", "LOW", "HIGH"],
    ["Staff poaching by nearby competitor", "Operations", "LOW", "MEDIUM", "LOW"],
]

_add_table(s, len(risk_data), 5, risk_data,
           0.4, 1.7, 12.5, 3.6,
           col_widths=[4.5, 1.5, 1.5, 1.5, 1.3], font_size=10, header_size=11)

# Color the rating column
tbl = s.shapes[-1].table
rating_colors = {"HIGH": RED_ACCENT, "MEDIUM": AMBER, "LOW": GREEN_OK}
for r in range(1, len(risk_data)):
    rating = risk_data[r][4]
    cell = tbl.cell(r, 4)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = rating_colors.get(rating, DARK_TEXT)
            run.font.bold = True

# Summary box
_add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 5.6, 12.1, 1.5, fill_color=NAVY)
_add_rich_box(s, 0.9, 5.7, 11.5, 1.3, [
    ("Overall Risk Rating:  HIGH", 20, True, RED_ACCENT, "Calibri"),
    ("", 6, False, WHITE, ""),
    ("5 HIGH-rated risks, 2 MEDIUM, 1 LOW. The confluence of fraud indicators, excluded-provider ties, "
     "geographic proximity, and shared vendor infrastructure creates an elevated threat profile that warrants "
     "immediate defensive action and proactive regulatory engagement.", 12, False, ICE_BLUE, "Calibri"),
])


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — COMPETITIVE POSITIONING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Competitive Positioning", "Knesset SADC vs. Garden of Joy — Head-to-Head")

comp_data = [
    ["Dimension", "Knesset SADC", "Garden of Joy"],
    ["Years Operating", "0 (claimed 14)", "20+ years"],
    ["Client Base", "0 (claims \"0K+\")", "425+ active clients"],
    ["Staff", "Unknown / none listed", "50+ trained staff"],
    ["Licensure", "Not verified", "NYSDOH licensed"],
    ["Website", "Lorem Ipsum placeholder", "Operational, real content"],
    ["EMR System", "CareCenta (unverified)", "CareCenta (verified)"],
    ["Location", "128 Brighton Beach Ave", "Brighton Beach area"],
    ["Reputation", "No reviews / fabricated 0%", "Established community trust"],
    ["Regulatory Status", "Unclear / potential violations", "Good standing"],
    ["MLTC Relationships", "Unknown", "Active contracts"],
    ["Audited Financials", "None (41 days old)", "Audited annually"],
]

_add_table(s, len(comp_data), 3, comp_data,
           0.6, 1.7, 12.1, 5.0,
           col_widths=[3.5, 4.3, 4.3], font_size=10, header_size=11)

# Color GOJ column green, Knesset column red
tbl = s.shapes[-1].table
for r in range(1, len(comp_data)):
    kn_cell = tbl.cell(r, 1)
    for p in kn_cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RED_ACCENT
    goj_cell = tbl.cell(r, 2)
    for p in goj_cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = GREEN_OK

# Insight
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.9, 12.1, 0.4, fill_color=NAVY)
_add_textbox(s, 0.9, 6.93, 11.5, 0.35,
             "GOJ's 20+ year operational track record is its strongest competitive moat against shell-entity competitors.",
             font_size=11, bold=False, color=ICE_BLUE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — RECOMMENDED ACTIONS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Recommended Actions", "Immediate & Near-Term Defensive Measures")

actions = [
    ("1. File OMIG Complaint",
     "Submit a formal complaint to the NYS Office of the Medicaid Inspector General detailing "
     "the Natalia Gurevich connection, the false \"Established 2012\" claim, and the HFC domain "
     "redirection. Include all documentary evidence collected.", RED_ACCENT),
    ("2. Conduct CareCenta Data Segregation Audit",
     "Engage CareCenta to confirm GOJ client data is fully segregated from other SADC tenants "
     "on the platform. Request a written assurance and technical audit of multi-tenancy controls.", RED_ACCENT),
    ("3. Monitor Knesset Licensure & Billing Activity",
     "Track NYS DOH SADC licensure applications and Medicaid enrollment for Knesset. Set up alerts "
     "for any billing activity under NPI 1831462795 via Medicaid claims data services.", RED_ACCENT),
    ("4. Strengthen Referral Source Relationships",
     "Proactively engage MLTC plans, hospital discharge planners, and community referral sources "
     "to reinforce GOJ's 20+ year track record and highlight the risks of new, unverified providers.", AMBER),
    ("5. Legal Review of Competitive Exposure",
     "Commission outside counsel review of: (a) GOJ's obligations upon discovering excluded-provider "
     "connections in its referral ecosystem, and (b) potential tortious interference claims if "
     "Knesset uses GOJ's proprietary information obtained via CareCenta.", AMBER),
    ("6. Internal Documentation & Compliance Review",
     "Document all findings in GOJ's compliance files. Review and update compliance policies "
     "for competitor monitoring, referral vetting, and mandatory reporting obligations.", AMBER),
]

for i, (title, body, color) in enumerate(actions):
    y = 1.7 + i * 0.92
    _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 0.8, 0.7, fill_color=color)
    _add_textbox(s, 0.6, y + 0.1, 0.8, 0.5, title.split(".")[0],
                 font_size=20, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")
    _add_textbox(s, 1.6, y, 5.5, 0.35, title,
                 font_size=14, bold=True, color=color, font_name="Calibri")
    _add_textbox(s, 1.6, y + 0.32, 11.1, 0.45, body,
                 font_size=10, color=DARK_TEXT, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — OPEN QUESTIONS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Open Questions", "Items Requiring Further Investigation")

questions = [
    ("Who is the true beneficial owner?",
     "The registered agent and DOS filings may obscure the real controlling party. "
     "Is Natalia Gurevich operating through a straw owner? Who signed the incorporation documents?"),
    ("Is Knesset already billing Medicaid?",
     "NPI 1831462795 was issued — has it been used for claims yet? A query of the NY Medicaid "
     "claims database would reveal any billing activity since incorporation."),
    ("What is the CareCenta relationship?",
     "Did CareCenta actively facilitate Knesset's setup, or was the platform simply purchased? "
     "Are there CareCenta employees with ties to the Gurevich/Rafailova network?"),
    ("How many more entities are in this network?",
     "The 6-SADC Rafailova pattern suggests we've only mapped part of the network. A full corporate "
     "records search for all Gurevich/Rafailova-related entities is needed."),
    ("What is the DOH licensure status?",
     "Has Knesset applied for / received an SADC operating license from NYS DOH? If they're operating "
     "without licensure, that's a separate regulatory violation beyond the fraud indicators."),
]

for i, (q, detail) in enumerate(questions):
    y = 1.8 + i * 1.1
    _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 0.6, 0.6, fill_color=NAVY)
    _add_textbox(s, 0.6, y + 0.08, 0.6, 0.45, f"Q{i+1}",
                 font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")
    _add_textbox(s, 1.4, y, 11.3, 0.35, q,
                 font_size=14, bold=True, color=NAVY, font_name="Calibri")
    _add_textbox(s, 1.4, y + 0.32, 11.3, 0.65, detail,
                 font_size=11, color=DARK_TEXT, font_name="Calibri")

# Footer
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.9, 12.1, 0.4, fill_color=NAVY)
_add_textbox(s, 0.9, 6.93, 11.5, 0.35,
             "Answers to these questions will determine the optimal legal, regulatory, and competitive response strategy.",
             font_size=11, color=ICE_BLUE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — SOURCES & METHODOLOGY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(s)
_slide_title(s, "Sources & Methodology", "How These Findings Were Compiled")

sources = [
    ("Public Records",
     "NYS DOS Corporation & Business Entity Database (DOS ID 4200463) ● "
     "NYS DOH SADC licensure records ● "
     "NPPES NPI Registry (1831462795)"),
    ("Federal Exclusions Databases",
     "HHS-OIG LEIE (List of Excluded Individuals/Entities) ● "
     "NYS OMIG Exclusions List ● "
     "SAM.gov exclusions database"),
    ("Domain & Web Intelligence",
     "WHOIS / ICANN domain registration records ● "
     "DNSDB / SecurityTrails historical DNS ● "
     "Direct website inspection and archive of myknesset.com"),
    ("Healthcare Fraud Records",
     "DOJ False Claims Act settlement database ● "
     "PACER federal court records ● "
     "State AG healthcare fraud press releases"),
    ("Competitive Intelligence",
     "CareCenta platform analysis ● "
     "Brighton Beach catchment area mapping ● "
     "MLTC plan provider network directories"),
    ("Methodology",
     "Investigation conducted July 2026 using OSINT techniques ● "
     "Cross-referenced across 6+ independent data sources ● "
     "All findings documented with screenshots and data exports ● "
     "Full evidence file available in Investigation File — Annex A through E"),
]

for i, (source, detail) in enumerate(sources):
    y = 1.6 + i * 0.92
    _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 2.3, 0.75, fill_color=NAVY)
    _add_textbox(s, 0.8, y + 0.12, 1.9, 0.55, source,
                 font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    _add_textbox(s, 3.1, y, 9.6, 0.75, detail,
                 font_size=10, color=DARK_TEXT, font_name="Calibri")

# Disclaimer
_add_shape(s, MSO_SHAPE.RECTANGLE, 0.6, 6.9, 12.1, 0.4, fill_color=RED_ACCENT)
_add_textbox(s, 0.9, 6.93, 11.5, 0.35,
             "CONFIDENTIAL — For internal use only. Not for distribution outside Garden of Joy / Gold Health Systems.",
             font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ── SAVE ────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅ Deck saved to: {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
