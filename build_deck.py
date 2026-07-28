#!/usr/bin/env python3
"""Build Kiselev-Gurevich SADC Network slide deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color scheme ──
DARK_BG     = RGBColor(0x0D, 0x0D, 0x0D)
RED_ACCENT  = RGBColor(0xCC, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY   = RGBColor(0x2A, 0x2A, 0x2A)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)
GREEN_FLAG  = RGBColor(0x00, 0x99, 0x00)
GOLD        = RGBColor(0xD4, 0xA0, 0x17)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=WHITE, spacing=1.2, font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return tf

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_classification_bar(slide):
    add_rect(slide, 0, 0, 13.333, 0.28, RED_ACCENT)
    tf = add_textbox(slide, 0.3, 0.02, 6, 0.25,
        "TOP SECRET // HCS // ORCON // NOFORN", font_size=8,
        color=WHITE, font_name="Courier New")
    add_textbox(slide, 10.3, 0.02, 3, 0.25,
        "JTF-BB-2026-0047-KG", font_size=8,
        color=WHITE, font_name="Courier New", alignment=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s, DARK_BG)
add_classification_bar(s)

# Red accent line
add_rect(s, 1.5, 2.1, 10.333, 0.04, RED_ACCENT)

add_textbox(s, 1.5, 2.3, 10.5, 1.2,
    "KISELEV-GUREVICH SADC NETWORK",
    font_size=42, bold=True, color=WHITE, font_name="Helvetica Neue")

add_textbox(s, 1.5, 3.5, 10.5, 0.8,
    "Joint Investigative Intelligence Briefing",
    font_size=24, color=LIGHT_GRAY)

add_multiline(s, 1.5, 4.6, 10.5, 1.5, [
    "Brooklyn, New York — Brighton Beach / Glen Head",
    "Medicaid Fraud Enterprise | False Claims Act | Enterprise Corruption",
    "Elderplan (31625) · Carecenta · HFC · Knesset · Garden of Joy",
    "16 July 2026",
], font_size=16, color=MID_GRAY, font_name="Courier New")

# Bottom bar
add_rect(s, 0, 7.2, 13.333, 0.3, RED_ACCENT)

# ════════════════════════════════════════════════════════
# SLIDE 2 — THE ENTERPRISE OVERVIEW
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "THE ENTERPRISE", font_size=32, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

add_multiline(s, 0.6, 1.5, 7.5, 3.5, [
    "$6.5M prior False Claims Act settlement (2018) — no CIA imposed",
    "Core entity: Home Family Care, Inc. (3051 Brighton 3rd St)",
    "421 members at Garden of Joy SADC — unregistered with DFTA",
    "Knesset SADC claims 'Established 2012' — incorporated June 2026",
    "Self-owned billing platform (Carecenta) controls all claims data",
    "Prior FCA defendants Kiselev + M. Gurevich both active in network",
    "Natalia Gurevich, MD: Enterprise Corruption felony; LEIE excluded",
    "Rafailova 6-sister patient steering network across 5+ SDACs",
], font_size=15, color=LIGHT_GRAY)

# Key stat boxes
for i, (num, label, color) in enumerate([
    ("$6.515M", "Prior FCA\nSettlement", RED_ACCENT),
    ("423", "GOJ Members\n(Unregistered)", RED_ACCENT),
    ("$15–20M/yr", "Est. Annual\nMedicaid Revenue", GOLD),
    ("$50–100M+", "Est. Total\nFCA Exposure", RED_ACCENT),
]):
    x = 8.8 + (i * 2.4) if i < 2 else 3.8 + ((i-2) * 4.5)
    y = 0.5 if i < 2 else 5.5
    rect = add_rect(s, x, y, 2.0, 1.5, DARK_GRAY, RED_ACCENT)
    add_textbox(s, x + 0.15, y + 0.15, 1.7, 0.6, num,
                font_size=26, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.15, y + 0.8, 1.7, 0.6, label,
                font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 3 — KEY PLAYERS
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "PRINCIPAL TARGETS", font_size=32, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

players = [
    ("ALEXANDER KISELEV", "UPenn Law '99, NY Bar\nPresident, HFC | Owner, Knesset SADC\n$6.415M FCA settlement (2018)\nGlen Head, NY — Uzbek immigrant",
     "PRIMARY\nTARGET", RED_ACCENT),
    ("MICHAEL GUREVICH", "VP, Home Family Care\nBoard-certified psychiatrist\n$100K FCA settlement (2018)\nNPI 1699773010 (active 7/2026)",
     "KEY\nWITNESS", GOLD),
    ("NATALIA GUREVICH", "Convicted: Enterprise Corruption (2004)\nLEIE excluded since 2005\nOMIG excluded — license revoked\nPossible mother of Michael Gurevich",
     "EXCLUDED\nFELON", RED_ACCENT),
    ("PETER ZESTYREV", "Chairman/CEO, Carecenta Inc.\nEx-CityMobil fintech (Moscow)\n~Zero public web presence\nUzbek surname — nominee CEO profile",
     "ENABLER", MID_GRAY),
    ("VLADIMIR KHIGER", "President, GOJ Inc.\n423 SADC members\nReportedly 'independent'\nOperational dependency on Carecenta",
     "STRAW\nOWNER", MID_GRAY),
    ("MARINA RABINOVICH", "1st Choice ADC + Direct Personal Care\nDFTA-registered SADC + home care\nSelf-referral / dual ownership\nMirrors HFC-Knesset-GOJ pattern",
     "PARALLEL\nTARGET", MID_GRAY),
]

for i, (name, details, badge, badge_color) in enumerate(players):
    row = i // 3
    col = i % 3
    x = 0.6 + (col * 4.2)
    y = 1.5 + (row * 2.85)

    # Card
    add_rect(s, x, y, 3.9, 2.6, DARK_GRAY, MID_GRAY)

    # Badge
    badge_box = add_rect(s, x + 0.15, y + 0.15, 0.9, 0.55, badge_color)
    add_textbox(s, x + 0.15, y + 0.17, 0.9, 0.5, badge,
                font_size=7, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
                font_name="Courier New")

    # Name
    add_textbox(s, x + 1.2, y + 0.15, 2.55, 0.5, name,
                font_size=11, bold=True, color=WHITE)

    # Details
    add_multiline(s, x + 0.15, y + 0.85, 3.6, 1.6, details.split("\n"),
                  font_size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════
# SLIDE 4 — ENTITY MAP
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "ENTITY RELATIONSHIP MAP", font_size=32, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

# Top: Kiselev
add_rect(s, 5.0, 1.4, 3.333, 0.6, RED_ACCENT)
add_textbox(s, 5.0, 1.45, 3.333, 0.5, "ALEXANDER KISELEV",
            font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Row 2: HFC, Knesset, Carecenta, GOJ
entities = [
    ("HFC Inc.\nLHCSA\n3051 Brighton 3rd St", 1.0, 2.6),
    ("Knesset SADC\nEst. ~2012 / Inc. 2026\n128 Brighton Beach Ave", 4.2, 2.6),
    ("Carecenta/Daycenta\nBilling Platform\n260 Madison Ave", 7.4, 2.6),
    ("GOJ Inc.\n423 members\nUNREGISTERED", 10.6, 2.6),
]
for label, x, y in entities:
    add_rect(s, x, y, 2.8, 1.0, DARK_GRAY, RED_ACCENT)
    add_textbox(s, x + 0.1, y + 0.08, 2.6, 0.9, label,
                font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Middle: Elderplan / Carecenta pipeline
add_rect(s, 4.0, 4.0, 5.333, 0.5, MID_GRAY)
add_textbox(s, 4.0, 4.05, 5.333, 0.4, "Elderplan (31625) ÷ Carecenta Billing Pipeline",
            font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom row: Gurevich + Rafailova
bottom = [
    ("Michael\nGurevich", "HFC VP\n$100K FCA", 1.8, 4.9),
    ("Natalia\nGurevich", "MD, Felon\nLEIE Excluded", 4.7, 4.9),
    ("Rafailova\n6 Sisters", "Patient Steering\nNetwork", 7.6, 4.9),
    ("Marina\nRabinovich", "1st Choice ADC\nSelf-Referral", 10.5, 4.9),
]
for name, sub, x, y in bottom:
    add_rect(s, x, y, 2.5, 0.9, DARK_GRAY, MID_GRAY)
    add_textbox(s, x + 0.1, y + 0.05, 2.3, 0.4, name,
                font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, y + 0.5, 2.3, 0.4, sub,
                font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Address nexus
add_multiline(s, 0.6, 6.2, 12, 1.0, [
    "ADDRESS NEXUS:  3051 Brighton 3rd St → HFC + Carecenta   |   128 Brighton Beach Ave #400A → Knesset + Carecenta   |   260 Madison Ave → Carecenta   |   Glen Head, NY → Kiselev + M. Gurevich + N. Rafailova",
], font_size=9, color=MID_GRAY, font_name="Courier New")

# ════════════════════════════════════════════════════════
# SLIDE 5 — SMOKING GUNS (part 1)
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "SMOKING GUNS", font_size=32, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

guns = [
    ("1", "$6.515M PRIOR FCA SETTLEMENT — NO CIA IMPOSED",
     "Kiselev and Michael Gurevich settled allegations of falsified attendance verification for $6.515M in 2018. HHS-OIG imposed zero Corporate Integrity Agreement — unprecedented for this magnitude. Both continue to operate and control SADC facilities billing the same Medicaid program through new entities.",
     RED_ACCENT),
    ("2", "EXCLUDED FELON PROVIDER IN THE NETWORK",
     "Natalia Gurevich, MD — convicted of Enterprise Corruption (NY's RICO equivalent, 2004), permanently excluded from ALL federal healthcare programs (LEIE + OMIG). Her continued proximity to billing entities means every claim submitted during any period of her involvement is a false claim.",
     RED_ACCENT),
    ("3", "KNESSET SADC: PHANTOM ENTITY — FALSIFIED 2012 ESTABLISHMENT DATE",
     "DFTA registration claims 'Established: 2012.' NY DOS shows Knesset Social Adult Day Care Inc. was incorporated June 2026. Website claims '15+ Years of Experience' for an entity that legally existed for weeks. DFTA registration filed under penalty of perjury.",
     RED_ACCENT),
]

for i, (num, title, detail, color) in enumerate(guns):
    y = 1.5 + (i * 1.95)
    # Number circle-ish
    add_rect(s, 0.6, y + 0.05, 0.6, 0.6, color)
    add_textbox(s, 0.6, y + 0.1, 0.6, 0.5, num,
                font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(s, 1.4, y, 11.5, 0.45, title, font_size=15, bold=True, color=WHITE)
    # Detail
    add_textbox(s, 1.4, y + 0.5, 11.5, 1.2, detail, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════
# SLIDE 6 — SMOKING GUNS (part 2)
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "DIGITAL FORENSICS & BILLING SMOKING GUNS", font_size=28, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

guns2 = [
    ("4", "DIGITAL FORENSIC TRAIL — DOMAIN SMOKING GUN",
     "Internet Archive proves myknesset.com redirected to hfcny.com through Feb 2018. WordPress admin email: tsigel@hfcny.com (Tatiana Sigel, HFC administrator). July 2025 website rebuild scrubbed all HFC references — spoliation concern. Knesset = rebranded HFC, not independent entity.",
     RED_ACCENT),
    ("5", "GARDEN OF JOY SADC — 423 MEMBERS, ZERO DFTA REGISTRATION",
     "GOJ Inc. (423 members, billing Elderplan via Carecenta) is NOT registered with NYC DFTA. NYC Admin. Code § 21-204 requires registration. Civil penalties: $250–$1,000/day. If operating 3+ years: potential $275K–$1M+ in accumulated fines.",
     RED_ACCENT),
    ("6", "CARECENTA — THE SELF-OWNED BILLING PLATFORM",
     "Carecenta's corporate address moved in lockstep with Kiselev entities: 3051 Brighton 3rd St → 128 Brighton Beach Ave → 260 Madison Ave. Platform operator has full competitor billing visibility. Only 2 named clients in 13 years — both Kiselev-controlled.",
     RED_ACCENT),
]

for i, (num, title, detail, color) in enumerate(guns2):
    y = 1.5 + (i * 1.95)
    add_rect(s, 0.6, y + 0.05, 0.6, 0.6, color)
    add_textbox(s, 0.6, y + 0.1, 0.6, 0.5, num,
                font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1.4, y, 11.5, 0.45, title, font_size=15, bold=True, color=WHITE)
    add_textbox(s, 1.4, y + 0.5, 11.5, 1.2, detail, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════
# SLIDE 7 — RECOMMENDED CHARGES
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "RECOMMENDED CHARGES", font_size=32, bold=True, color=WHITE)
add_rect(s, 0.6, 1.15, 2.5, 0.04, RED_ACCENT)

charges = [
    ("18 U.S.C. § 1347\nHealthcare Fraud", "Kiselev,\nM. Gurevich", "Ongoing false claims to\nElderplan via Carecenta.\nPrior FCA pattern.", "10–20 yrs\nper count"),
    ("31 U.S.C. § 3729\nFalse Claims Act", "HFC, Knesset,\nGOJ, Kiselev", "N. Gurevich involvement +\nfalse DFTA registration.\nEvery claim a false claim.", "Treble damages\n+ $13.5K–$27K/claim"),
    ("18 U.S.C. § 1956\nMoney Laundering", "Kiselev,\nZestyrev", "Carecenta billing pipeline;\nproperty holdings;\nUzbek transnational.", "20 yrs +\n$500K or 2×"),
    ("18 U.S.C. § 1962\nRICO", "Kiselev,\nM. Gurevich,\nN. Gurevich", "Multi-entity enterprise;\nMedicaid fraud pattern\nacross 2 generations.", "20 yrs +\nforfeiture"),
    ("42 U.S.C. § 1320a-7\nExclusion Violation", "All entities\nw/ N. Gurevich", "Payment to excluded\nindividual; OIG LEIE\nscreening failure.", "CMP up to\n$20K/item"),
    ("18 U.S.C. § 1001\nFalse Statements", "Kiselev", "DFTA 'Est. 2012' for\nentity incorporated\nJune 2026.", "5 yrs\nper count"),
]

for i, (statute, target, basis, penalty) in enumerate(charges):
    row = i // 3
    col = i % 3
    x = 0.6 + (col * 4.2)
    y = 1.5 + (row * 2.85)

    add_rect(s, x, y, 3.9, 2.6, DARK_GRAY, MID_GRAY)

    # Statute header
    add_rect(s, x, y, 3.9, 0.6, DARK_BG)
    add_textbox(s, x + 0.15, y + 0.04, 3.6, 0.55, statute,
                font_size=9, bold=True, color=RED_ACCENT, font_name="Courier New")

    # Target
    add_textbox(s, x + 0.15, y + 0.7, 1.5, 0.4, "TARGET:", font_size=7, color=MID_GRAY)
    add_textbox(s, x + 0.15, y + 0.9, 3.6, 0.6, target,
                font_size=9, bold=True, color=WHITE)

    # Basis
    add_textbox(s, x + 0.15, y + 1.5, 1.5, 0.3, "BASIS:", font_size=7, color=MID_GRAY)
    add_textbox(s, x + 0.15, y + 1.65, 3.6, 0.7, basis,
                font_size=8, color=LIGHT_GRAY)

    # Penalty
    add_rect(s, x, y + 2.2, 3.9, 0.4, RED_ACCENT)
    add_textbox(s, x + 0.15, y + 2.22, 3.6, 0.35, penalty,
                font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
                font_name="Courier New")

# ════════════════════════════════════════════════════════
# SLIDE 8 — IMMEDIATE ACTIONS
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_classification_bar(s)

add_textbox(s, 0.6, 0.5, 12, 0.6, "IMMEDIATE INVESTIGATIVE ACTIONS", font_size=28, bold=True, color=WHITE)
add_rect(s, 0.6, 1.05, 2.5, 0.04, RED_ACCENT)

actions = [
    ("Evidentiary Subpoenas", [
        "→ Elderplan (31625): all claims for Knesset, GOJ, HFC — pattern analysis",
        "→ Carecenta: audit logs, source code, raw claims — preserve NOW",
        "→ NY DFS: property records, corporate filings, financial accounts",
    ]),
    ("Regulatory Referrals", [
        "→ DFTA: GOJ unregistered operation + Knesset phantom establishment date",
        "→ HHS-OIG / OMIG: Natalia Gurevich exclusion violation investigation",
        "→ NY AG MFCU: multi-agency task force — prior $6.5M settlement",
        "→ NY Bar: Kiselev Rule 8.4 (dishonesty, fraud, misrepresentation)",
    ]),
    ("Field Operations", [
        "→ Undercover test visits — GOJ & Knesset facilities",
        "→ Financial investigation — trace settlement source + asset shielding",
        "→ CBP/ICE liaison — Uzbek transnational money movement",
        "→ Rafailova network — patient steering + referral remuneration",
    ]),
]

left_cols = [0.6, 4.8, 9.0]
for i, (title, items) in enumerate(actions):
    x = left_cols[i]
    add_rect(s, x, 1.4, 3.9, 0.5, DARK_GRAY)
    add_textbox(s, x + 0.15, 1.43, 3.6, 0.45, title,
                font_size=14, bold=True, color=RED_ACCENT, font_name="Courier New")
    add_multiline(s, x + 0.15, 2.05, 3.7, 4.5, items,
                  font_size=11, color=LIGHT_GRAY, spacing=1.4)

# Bottom assessment
add_rect(s, 0, 7.0, 13.333, 0.5, RED_ACCENT)
add_textbox(s, 0.6, 7.03, 12, 0.45,
    "ASSESSMENT: Active, ongoing multi-million dollar Medicaid fraud enterprise. Prior $6.5M FCA settlement + excluded felon provisioner + self-owned billing platform + unregistered 423-member SADC + transnational dimension. IMMEDIATE MULTI-AGENCY INVESTIGATION WITH SUBPOENA POWER WARRANTED.",
    font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ── Save ──
output = "/Users/mainsobhelper/Desktop/REX/Kiselev-Gurevich-Network-Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
